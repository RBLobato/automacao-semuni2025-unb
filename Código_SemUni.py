# ===== Importar as bibliotecas =====

import pandas as pd 
from pptx import Presentation 
from pptx.util import Inches, Pt, Cm, Mm
from pptx.enum.text import PP_ALIGN 
from pptx.dml.color import RGBColor
from math import ceil

# ===== Planilha do Excel =====

planilha = "planilha_semuni_limpa.xlsx"

#===== Nomes das Colunas na Planilha =====

COLS = {
"título" : "Título do projeto" , 
"coordenador" : "Coordenador" ,
"participantes" : "Participantes" ,
"resumo" : "Resumo do projeto" ,
"palavras-chave" : "Palavras-Chaves" ,
"area_principal" : "Área" ,
}

# ===== Cores, fontes e tamanhos =====

roxo_hex = "#922F60"   
bege_hex = "#E2D2AF"   
fonte = "Aptos (Corpo)"     
tam_titulo = 14
tam_text = 12
borda_pts = 3

# ===== Leitura da planilha =====

df = pd.read_excel(planilha)

# ===== Verificação da existência das colunas =====

faltando = [COLS[k] for k in COLS if COLS[k] not in df.columns]
if faltando:
    raise ValueError(f"As seguintes colunas não foram encontradas no Excel: {faltando}")

# ===== Criando a apresentação PowerPoint =====

prs = Presentation()

prs.slide_width  = Cm(21.0) # largura A4 em mm
prs.slide_height = Cm(29.7) # altura A4 em mm

# prs.save("projetos_A4_inicial.pptx") # Criação da planilha de teste.

print("OK: planilha lida e apresentação A4 criada (projetos_A4_inicial.pptx).")

#===== Converter HEX para RGB =====

def hex2rgb(cor_hex: str):
    h = cor_hex.strip().lstrip("#")
    return tuple(int(h[i:i+2] , 16) for i in (0 , 2 , 4))

RGB_ROXO = RGBColor(*hex2rgb(roxo_hex))   
RGB_BEGE = RGBColor(*hex2rgb(bege_hex))   

def add_textbox(
    slide,
    left_cm, top_cm, width_cm, height_cm,
    fill_rgb=None, border_rgb=None, border_pts=None
):
    shape = slide.shapes.add_textbox(Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm))

     # Preenchimento
    if fill_rgb is not None:
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    # Borda
    if border_rgb is not None and border_pts is not None:
        shape.line.color.rgb = border_rgb
        shape.line.width = Pt(border_pts)
    else:
        shape.line.fill.background()
    # Quebra automática
    shape.text_frame.word_wrap = True
    return shape

def add_label_value(shape, label_text, value_text, tamanho_pt, cor_rgb, fonte_nome, alinhamento=PP_ALIGN.JUSTIFY):
    """Escreve 'Label: Valor' na mesma caixa (label em negrito, valor normal)."""
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = alinhamento

    r1 = p.add_run()
    r1.text = label_text
    r1.font.bold = True
    r1.font.size = Pt(tamanho_pt)
    r1.font.color.rgb = cor_rgb
    r1.font.name = fonte_nome

    r2 = p.add_run()
    r2.text = f" {value_text}".strip()
    r2.font.bold = False
    r2.font.size = Pt(tamanho_pt)
    r2.font.color.rgb = cor_rgb
    r2.font.name = fonte_nome

from math import ceil

def _chars_por_linha_aprox(largura_cm: float, font_pt: float) -> float:
    """
    Estima quantos caracteres cabem por linha para a largura e fonte dadas.
    Heurística: largura em pontos / (0.5 * font_pt). Ajuste o fator 0.5 se quiser.
    """
    largura_in = largura_cm / 2.54
    largura_pt = largura_in * 72.0
    fator_largura_char = 0.5  # <- ajuste fino: 0.45 (mais chars por linha) / 0.55 (menos)
    return max(8.0, largura_pt / (fator_largura_char * font_pt))

def _linhas_estimadas(texto: str, largura_cm: float, font_pt: float) -> int:
    """
    Estima o total de linhas levando em conta quebras de parágrafo (\n).
    Usa uma conta aproximada por parágrafo baseada em caracteres/linha.
    """
    cpl = _chars_por_linha_aprox(largura_cm, font_pt)
    total_linhas = 0
    for par in (texto or "").splitlines() or [""]:
        # pequeno bônus para espaços/pontuação que “quebram” melhor
        fator_respiro = 0.95
        n_chars = len(par.strip())
        linhas = 1 if n_chars == 0 else ceil((n_chars * fator_respiro) / cpl)
        total_linhas += max(1, linhas)
    return total_linhas

def altura_resumo_cm(texto: str, largura_cm: float, font_pt: float,
                     min_cm: float = 6.0, max_cm: float | None = None) -> float:
    """
    Converte linhas estimadas -> altura da caixa em cm.
    Usa line-height ≈ 1.2 * font_pt. Adiciona padding superior/inferior.
    Respeita limites min/max em cm.
    """
    linhas = _linhas_estimadas(texto, largura_cm, font_pt)
    line_height_pt = 1.2 * font_pt
    # converte points -> cm (1 pt = 1/72 in; 1 in = 2.54 cm)
    line_height_cm = line_height_pt * (2.54 / 72.0)
    padding_cm = 0.6  # espaço de respiro dentro da caixa (ajustável)
    h_cm = linhas * line_height_cm + padding_cm

    if max_cm is not None:
        h_cm = min(h_cm, max_cm)
    h_cm = max(h_cm, min_cm)
    return h_cm


#===== Layout dos Slides =====

MARGEM_L = 2.0
MARGEM_R = 2.0
LARGURA_TOTAL = 21.0  # A4
ALTURA_TOTAL  = 29.7  # A4
BOX_W = LARGURA_TOTAL - (MARGEM_L + MARGEM_R)

ESPACO_V = 0.4  # espaço vertical entre caixas

ALT_TITULO = 3
ALT_COORD  = 1
ALT_PART   = 1
ALT_PCHAVE = 1
ALT_AREA   = 1
ALT_RESUMO = 2.2 
MARGEM_SUP = 1.5
MARGEM_INF = 1.5

top = MARGEM_SUP

#===== Criação dos Slides =====

df_f = df.fillna("")  # evita aparecer "nan"

for _, row in df_f.iterrows():
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # slide em branco
    top = MARGEM_SUP  # margem superior ajustável

    # ===== 1) TÍTULO (altura automática) =====
    texto_titulo = str(row[COLS["título"]]).strip()
    alt_titulo_auto = altura_resumo_cm(
        texto=texto_titulo,
        largura_cm=BOX_W,
        font_pt=tam_titulo,
        min_cm=1.8,   # altura mínima
        max_cm=3.5    # altura máxima (caso o título seja muito longo)
    )

    titulo_shape = add_textbox(
        slide,
        left_cm=MARGEM_L, top_cm=top, width_cm=BOX_W, height_cm=alt_titulo_auto,
        fill_rgb=RGB_ROXO, border_rgb=None, border_pts=None
    )
    tf = titulo_shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = texto_titulo
    r.font.bold = True
    r.font.size = Pt(tam_titulo)
    r.font.color.rgb = RGB_BEGE
    r.font.name = fonte

    top += alt_titulo_auto + ESPACO_V

    # ===== 2) COORDENADOR(A) (altura automática) =====
    texto_coord = str(row[COLS["coordenador"]]).strip()
    alt_coord_auto = altura_resumo_cm(
        texto=texto_coord,
        largura_cm=BOX_W,
        font_pt=tam_text,
        min_cm=1.8,
        max_cm=3.0
    )

    coord_shape = add_textbox(
        slide,
        left_cm=MARGEM_L, top_cm=top, width_cm=BOX_W, height_cm=alt_coord_auto,
        fill_rgb=RGB_BEGE, border_rgb=RGB_ROXO, border_pts=borda_pts
    )
    add_label_value(
        coord_shape, "Coordenador(a):", texto_coord,
        tam_text, RGB_ROXO, fonte, alinhamento=PP_ALIGN.JUSTIFY
    )

    top += alt_coord_auto + ESPACO_V

    # ===== 3) PARTICIPANTES (altura automática) =====
    texto_part = str(row[COLS["participantes"]]).strip()
    alt_part_auto = altura_resumo_cm(
        texto=texto_part,
        largura_cm=BOX_W,
        font_pt=tam_text,
        min_cm=2.0,
        max_cm=8.0
    )

    part_shape = add_textbox(
        slide,
        left_cm=MARGEM_L, top_cm=top, width_cm=BOX_W, height_cm=alt_part_auto,
        fill_rgb=RGB_BEGE, border_rgb=RGB_ROXO, border_pts=borda_pts
    )
    add_label_value(
        part_shape, "Participantes:", texto_part,
        tam_text, RGB_ROXO, fonte, alinhamento=PP_ALIGN.JUSTIFY
    )

    top += alt_part_auto + ESPACO_V

    # ===== 4) PALAVRAS-CHAVE (altura automática) =====
    texto_pch = str(row[COLS["palavras-chave"]]).strip()
    alt_pch_auto = altura_resumo_cm(
        texto=texto_pch,
        largura_cm=BOX_W,
        font_pt=tam_text,
        min_cm=2.0,
        max_cm=6.0
    )

    pch_shape = add_textbox(
        slide,
        left_cm=MARGEM_L, top_cm=top, width_cm=BOX_W, height_cm=alt_pch_auto,
        fill_rgb=RGB_BEGE, border_rgb=RGB_ROXO, border_pts=borda_pts
    )
    add_label_value(
        pch_shape, "Palavras-Chave:", texto_pch,
        tam_text, RGB_ROXO, fonte, alinhamento=PP_ALIGN.JUSTIFY
    )

    top += alt_pch_auto + ESPACO_V

    # ===== 5) ÁREA PRINCIPAL (fixa — geralmente curta) =====
    texto_area = str(row[COLS["area_principal"]]).strip()
    alt_area_auto = altura_resumo_cm(
        texto=texto_area,
        largura_cm=BOX_W,
        font_pt=tam_text,
        min_cm=2.0,
        max_cm=3.0
    )

    area_shape = add_textbox(
        slide,
        left_cm=MARGEM_L, top_cm=top, width_cm=BOX_W, height_cm=alt_area_auto,
        fill_rgb=RGB_BEGE, border_rgb=RGB_ROXO, border_pts=borda_pts
    )
    add_label_value(
        area_shape, "Área Principal:", texto_area,
        tam_text, RGB_ROXO, fonte, alinhamento=PP_ALIGN.JUSTIFY
    )

    top += alt_area_auto + ESPACO_V

    # ===== 6) RESUMO (altura automática) =====
    texto_resumo = str(row[COLS["resumo"]]).strip()
    max_alt_resumo = max(3.0, ALTURA_TOTAL - MARGEM_INF - top)
    alt_resumo_auto = altura_resumo_cm(
        texto=texto_resumo,
        largura_cm=BOX_W,
        font_pt=tam_text,
        min_cm=6.0,
        max_cm=max_alt_resumo
    )

    resumo_shape = add_textbox(
        slide,
        left_cm=MARGEM_L, top_cm=top, width_cm=BOX_W, height_cm=alt_resumo_auto,
        fill_rgb=RGB_BEGE, border_rgb=RGB_ROXO, border_pts=borda_pts
    )
    tf = resumo_shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.JUSTIFY
    r = p.add_run()
    r.text = texto_resumo
    r.font.bold = False
    r.font.size = Pt(tam_text)
    r.font.color.rgb = RGB_ROXO
    r.font.name = fonte


# ===== Salvar =====
prs.save("projetos_A4_etapa2.pptx")
print("✅ OK: slides gerados em 'projetos_A4_etapa2.pptx'.")
